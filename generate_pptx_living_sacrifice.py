"""Generate PPTX for: ECC Sunday School - Becoming a Living Sacrifice (Slides)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Palette --
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
BG_INDIGO    = RGBColor(0x1E, 0x1B, 0x4B)
BG_TEAL      = RGBColor(0x13, 0x4E, 0x4A)
BG_BLUE      = RGBColor(0x0C, 0x4A, 0x6E)
BG_PURPLE    = RGBColor(0x4A, 0x19, 0x42)
BG_GREEN     = RGBColor(0x14, 0x53, 0x2D)
BG_STONE     = RGBColor(0x1C, 0x19, 0x17)

WHITE        = RGBColor(0xF1, 0xF5, 0xF9)
GOLD         = RGBColor(0xFB, 0xBF, 0x24)
LIGHT_BLUE   = RGBColor(0x93, 0xC5, 0xFD)
INDIGO_ACC   = RGBColor(0x81, 0x8C, 0xF8)
TEAL_ACC     = RGBColor(0x5E, 0xEA, 0xD4)
PURPLE_ACC   = RGBColor(0xF0, 0xAB, 0xFC)
GREEN_ACC    = RGBColor(0x86, 0xEF, 0xAC)
RED_ACC      = RGBColor(0xF8, 0x71, 0x71)
LIGHT_RED    = RGBColor(0xFC, 0xA5, 0xA5)
LIGHT_GREEN  = RGBColor(0xBB, 0xF7, 0xD0)
SLATE        = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MAIN    = RGBColor(0xE2, 0xE8, 0xF0)
CYAN_ACC     = RGBColor(0x7D, 0xD3, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=32, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf


def add_para(tf, text, font_size=32, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    return p


# -- Build --
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# -- SLIDE 1: Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "ECC REDMOND  |  MARCH 22, 2026", Inches(1), Inches(1.8),
         Inches(11.3), Inches(0.6), font_size=20, bold=True,
         color=INDIGO_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Becoming a Living Sacrifice", Inches(1), Inches(2.5),
         Inches(11.3), Inches(1.4), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Sunday School \u2014 Middle Schoolers (Grades 6\u20138)", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, "Romans 12:1-8", Inches(1), Inches(4.9),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)

# -- SLIDE 2: Icebreaker Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_TEAL)
add_text(s, "\U0001f3b0", Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Lottery Life", Inches(1), Inches(2.0),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "A poll from Philadelphia asked: You just won $100 MILLION.\n"
         "You never have to work again. No homework. No chores.\n"
         "Nobody can tell you what to do.",
         Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.4),
         font_size=26, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "Total freedom.", Inches(1), Inches(5.0),
         Inches(11.3), Inches(0.7), font_size=36, bold=True,
         color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s, "What are the TOP 3 things you'd do with your life?", Inches(1), Inches(5.7),
         Inches(11.3), Inches(0.7), font_size=32, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "You have 30 seconds. GO!", Inches(1), Inches(6.5),
         Inches(11.3), Inches(0.5), font_size=22, italic=True,
         color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 3: Reveal intro --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x0D, 0x94, 0x88))
add_text(s, "\U0001f4ca", Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s,
         "Let's see how your answers stack up\nagainst a REAL POLL from Philadelphia!",
         Inches(1), Inches(3.0), Inches(11.3), Inches(1.6),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# -- SLIDES 4-8: Lottery answers --
lottery_data = [
    ("Poll Answer #1", "\U0001f30d", "Travel the world"),
    ("Poll Answer #2", "\U0001f4b3", "Pay off all debts"),
    ("Poll Answer #3", "\U0001f3a8", "Become an artist"),
    ("Poll Answer #4", "\U0001f4d6", "Write a book"),
    ("Poll Answer #5", "\U0001f3ae", "Play video games full time"),
]
subtitles = [
    None, None, None, None,
    'Pro gamer status \u2014 Elder Elton said\n"this might ring true to some of the youth here"',
]

for (label, emoji, text), sub in zip(lottery_data, subtitles):
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_TEAL)
    add_text(s, label, Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
             font_size=22, bold=True, color=TEAL_ACC, align=PP_ALIGN.CENTER)
    add_text(s, emoji, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
             font_size=72, align=PP_ALIGN.CENTER)
    add_text(s, text, Inches(1), Inches(3.3), Inches(11.3), Inches(1.0),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, sub, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.9),
                 font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 9: Transition — Spiritual Lottery --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "We've won the SPIRITUAL lottery.", Inches(1), Inches(0.5), Inches(11.3), Inches(1.0),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "Jesus paid off ALL our debts. Last week, Minister Ryan\n"
         "preached about our freedom from condemnation (Romans 8).",
         Inches(1.2), Inches(1.7), Inches(10.8), Inches(1.2),
         font_size=28, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         'A lot of Christians think: "Once I\'m saved, that\'s it!\n'
         'I\'m on Easy Street until heaven."',
         Inches(1.2), Inches(3.2), Inches(10.8), Inches(1.0),
         font_size=28, color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s,
         "But there are 16 chapters in Romans. If salvation was\n"
         "all there was, it would end at chapter 11.",
         Inches(1.2), Inches(4.4), Inches(10.8), Inches(1.0),
         font_size=28, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "So what are you going to do\nwith your freedom?",
         Inches(1), Inches(5.8), Inches(11.3), Inches(1.0),
         font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 10: Scripture Romans 12:1 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Key Verse", Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"Therefore, I urge you, brothers and sisters,\n'
         'in view of God\u2019s mercy, to offer your bodies\n'
         'as a living sacrifice, holy and pleasing to God \u2014\n'
         'this is your true and proper worship."',
         Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.8),
         font_size=32, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Romans 12:1", Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 11: Scripture Romans 12:2 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s,
         '"Do not conform to the pattern of this world,\n'
         'but be transformed by the renewing of your mind.\n'
         'Then you will be able to test and approve what\n'
         'God\u2019s will is \u2014 his good, pleasing and perfect will."',
         Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.8),
         font_size=32, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Romans 12:2", Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 12: Sermon Outline --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "Last Week's Sermon", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=INDIGO_ACC, align=PP_ALIGN.LEFT)
add_text(s, "Becoming a Living Sacrifice", Inches(1), Inches(1.5), Inches(11.3), Inches(0.9),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
parts = [
    "Part 1:  What Is a Living Sacrifice?",
    "Part 2:  Three Choices to Surrender",
    "Part 3:  True Worship = Service",
]
y = Inches(3.0)
for part in parts:
    add_text(s, "\u25b8  {}".format(part), Inches(1), y, Inches(11), Inches(0.7),
             font_size=32, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 13: Key Idea 1 — Oxymoron --
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE)
add_text(s, "Key Idea 1", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=LIGHT_BLUE)
add_text(s, "What Is a Living Sacrifice?", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
add_text(s, '"Living sacrifice" is an OXYMORON!', Inches(0.8), Inches(2.0),
         Inches(11.5), Inches(0.6), font_size=30, bold=True, color=GOLD)
items = [
    "In the OT, when you sacrifice something, it's DEAD. Pigeons dead. Rams dead. Goats dead.",
    "Living = alive in God because of Christ (faith, not works!)",
    "Sacrifice = giving yourself FULLY to God \u2014 no compartmentalizing",
    "NOT reintroducing OT sacrifices. Jesus already paid for sins \u2014 that is DONE.",
]
y = Inches(2.8)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 14: Keller Quote --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s,
         '"To live a Christian life, you\u2019re not living\n'
         'a Christian life, unless you put to death\n'
         'the idea that you have a right to live\n'
         'as you choose."',
         Inches(1.2), Inches(1.5), Inches(10.8), Inches(3.5),
         font_size=34, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Timothy Keller", Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 15: Key Idea 2 Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_PURPLE)
add_text(s, "Key Idea 2", Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=PURPLE_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Three Choices to Surrender", Inches(0.8), Inches(2.3),
         Inches(11.5), Inches(1.0), font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Becoming a living sacrifice means handing over\nthree big choices to God...",
         Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.2),
         font_size=28, color=TEXT_MAIN, align=PP_ALIGN.CENTER)

# -- SLIDE 16: Surrender 1 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x58, 0x1C, 0x87))
add_text(s, "1", Inches(1), Inches(0.5), Inches(11.3), Inches(1.2),
         font_size=72, bold=True, color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s, "I'm letting God choose what is\nRIGHT AND WRONG",
         Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         font_size=38, bold=True, color=PURPLE_ACC, align=PP_ALIGN.CENTER)
add_text(s, '"Do not conform to the pattern of this world"',
         Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.6),
         font_size=26, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
items = [
    "If God told you that you couldn't lie anymore \u2014 would that cramp your lifestyle?",
    "Example: parents lying about their kid's age for a discount",
    "As Christians: we trust God's wisdom to define right and wrong",
]
y = Inches(4.2)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.8)

# -- SLIDE 17: Surrender 2 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x7C, 0x2D, 0x8E))
add_text(s, "2", Inches(1), Inches(0.5), Inches(11.3), Inches(1.2),
         font_size=72, bold=True, color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s, "I'm letting God choose\nWHO I'M BECOMING",
         Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         font_size=38, bold=True, color=PURPLE_ACC, align=PP_ALIGN.CENTER)
add_text(s, '"Be transformed by the renewing of your mind"',
         Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.6),
         font_size=26, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
items = [
    "God changes you from the inside out \u2014 character, heart, perspective",
    "Transformation = becoming more like Christ",
    "The CBE study shows what happens when you let God's Word transform you...",
]
y = Inches(4.3)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.8)

# -- SLIDE 18: Bible Reading Stats --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x03, 0x69, 0xA1))
add_text(s, "Bible Reading 4+ Times/Week", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.6), font_size=28, bold=True,
         color=CYAN_ACC, align=PP_ALIGN.CENTER)
add_text(s, "CBE Study of 40,000 People", Inches(1), Inches(1.2),
         Inches(11.3), Inches(0.6), font_size=34, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

stats = [
    ("-30%", "Loneliness"),
    ("-32%", "Anger Issues"),
    ("-40%", "Bitterness in Relationships"),
    ("-61%", "Viewing Pornography"),
    ("-74%", "Gambling"),
    ("+407%", "Memorizing Scripture"),
    ("+228%", "Sharing Your Faith"),
]
# Layout: 4 in left column, 3 in right column
col1_x = Inches(1.5)
col2_x = Inches(7.0)
y_start = Inches(2.2)
for i, (num, label) in enumerate(stats):
    if i < 4:
        x = col1_x
        y = y_start + Inches(i * 1.2)
    else:
        x = col2_x
        y = y_start + Inches((i - 4) * 1.2)
    add_text(s, num, x, y, Inches(4.5), Inches(0.7),
             font_size=38, bold=True, color=GOLD)
    add_text(s, label, x, y + Inches(0.5), Inches(4.5), Inches(0.5),
             font_size=20, color=SLATE)

# -- SLIDE 19: Surrender 3 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x93, 0x33, 0xA0))
add_text(s, "3", Inches(1), Inches(0.5), Inches(11.3), Inches(1.2),
         font_size=72, bold=True, color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s, "I'm letting God choose\nWHAT I'LL DO WITH MY LIFE",
         Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         font_size=38, bold=True, color=PURPLE_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Romans 12:3-8 \u2014 Using your gifts for others, not just yourself",
         Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.6),
         font_size=26, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
items = [
    'The opposite = "quiet quitting" on God \u2014 just doing the minimum',
    '"Whatever you do, work at it with all your heart,\n   as working for the Lord" \u2014 Colossians 3',
    "Your talents and gifts are for the benefit of OTHERS",
]
y = Inches(4.3)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.8),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 20: Three Surrenders Summary --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x29, 0x25, 0x24))
add_text(s, "Three Surrenders", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.8), font_size=34, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

surrenders = [
    "\U0001f9ed  What's right & wrong",
    "\U0001f9e0  Who I'm becoming",
    "\U0001f3af  What I'll do with my life",
]
y = Inches(2.0)
for v in surrenders:
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(3), y, Inches(7.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.fill.fore_color.brightness = -0.85
    shape.line.color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].text = v
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    y += Inches(1.3)

add_text(s, "The Surrendered Life = giving up your right to choose\nyour own way and following God's way instead.",
         Inches(1), Inches(6.0), Inches(11.3), Inches(1.0),
         font_size=24, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 21: Key Idea 3 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "Key Idea 3", Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=GREEN_ACC)
add_text(s, "True Worship = Service", Inches(0.8), Inches(1.2),
         Inches(11.5), Inches(0.9), font_size=46, bold=True, color=WHITE)
items = [
    'Greek word for "worship" in Romans 12:1 = latreia',
    "It means BOTH worship AND service",
    "True worship isn't just singing on Sunday \u2014 it's how you live Mon\u2013Sat",
    "The motivation isn't GUILT \u2014 it's GRATITUDE",
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 22: Selfishness vs Selflessness --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x29, 0x25, 0x24))
add_text(s, "The Big Shift", Inches(1), Inches(0.4),
         Inches(11.3), Inches(0.8), font_size=34, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

# Left column
add_text(s, "Living for Yourself", Inches(0.5), Inches(1.5),
         Inches(5.8), Inches(0.6), font_size=24, bold=True, color=RED_ACC,
         align=PP_ALIGN.CENTER)
self_items = [
    "Use gifts to promote yourself",
    "Maximize your own comfort",
    '"Quiet quitting" \u2014 do the minimum',
    "Self-centered perspective",
]
y = Inches(2.3)
for item in self_items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(5.2), Inches(0.6),
             font_size=24, color=LIGHT_RED)
    y += Inches(0.7)

# Right column
add_text(s, "Living as a Sacrifice", Inches(6.8), Inches(1.5),
         Inches(5.8), Inches(0.6), font_size=24, bold=True, color=GREEN_ACC,
         align=PP_ALIGN.CENTER)
sacrifice_items = [
    "Use gifts to serve others",
    "Put others' needs first",
    "Work with all your heart for the Lord",
    "Selfless, Christ-like love",
]
y = Inches(2.3)
for item in sacrifice_items:
    add_text(s, "\u25b8  {}".format(item), Inches(7.1), y, Inches(5.2), Inches(0.6),
             font_size=24, color=LIGHT_GREEN)
    y += Inches(0.7)

# -- SLIDE 23: Video --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x37, 0x30, 0xA3))
add_text(s, "Recommended Video (8 min)", Inches(1), Inches(1.5),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
add_text(s, "BibleProject: Romans 5\u201316 Overview", Inches(1), Inches(2.3),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "\u25b6  youtube.com/watch?v=0SVTl4Xa5fY", Inches(1), Inches(3.8),
         Inches(11.3), Inches(0.6), font_size=28, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         'Click link above or search "BibleProject Romans 5-16"\n'
         "Covers Romans 5\u201316 themes including living sacrifice and the body of Christ",
         Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 24: Key Takeaway --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "THE BOTTOM LINE", Inches(1), Inches(1.0),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Jesus paid off all your debts.\nYou are FREE.",
         Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
         font_size=40, italic=True, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "The question is: what are you going to do\nwith that freedom?",
         Inches(1), Inches(3.8), Inches(11.3), Inches(1.0),
         font_size=34, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "Not guilt. GRATITUDE.", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.8), font_size=46, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 25: Small Group Time --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "\U0001f4ac", Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Small Group Time", Inches(1), Inches(2.8),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Split into groups of 4-5 with a leader", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s,
         "No judgment zone. Share honestly.\nAsk questions. Listen to each other.",
         Inches(1), Inches(5.0), Inches(11.3), Inches(1.0),
         font_size=24, color=INDIGO_ACC, align=PP_ALIGN.CENTER)

# -- Save --
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ECC Sunday School - Becoming a Living Sacrifice (Slides).pptx")
prs.save(out)
print("Saved -> {}".format(out))
print("Total slides: {}".format(len(prs.slides)))
