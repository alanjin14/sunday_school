"""Generate PPTX for: ECC Sunday School - The Power to Set Free (Slides)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
BG_INDIGO    = RGBColor(0x1E, 0x1B, 0x4B)
BG_TEAL      = RGBColor(0x13, 0x4E, 0x4A)
BG_BLUE      = RGBColor(0x0C, 0x4A, 0x6E)
BG_PURPLE    = RGBColor(0x4A, 0x19, 0x42)
BG_GREEN     = RGBColor(0x14, 0x53, 0x2D)
BG_STONE     = RGBColor(0x1C, 0x19, 0x17)

WHITE        = RGBColor(0xF1, 0xF5, 0xF9)
GOLD         = RGBColor(0xFB, 0xBF, 0x24)
LIGHT_BLUE   = RGBColor(0x93, 0xC5, 0xFD)
INDIGO_ACC   = RGBColor(0x81, 0x8C, 0xF8)
TEAL_ACC     = RGBColor(0x5E, 0xEA, 0xD4)
PURPLE_ACC   = RGBColor(0xF0, 0xAB, 0xFC)
GREEN_ACC    = RGBColor(0x86, 0xEF, 0xAC)
RED_ACC      = RGBColor(0xF8, 0x71, 0x71)
LIGHT_RED    = RGBColor(0xFC, 0xA5, 0xA5)
LIGHT_GREEN  = RGBColor(0xBB, 0xF7, 0xD0)
SLATE        = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MAIN    = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=32, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf


def add_para(tf, text, font_size=32, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    return p


def bullet_slide(slide, title, items, bg_color=BG_BLUE,
                 title_color=LIGHT_BLUE, item_color=TEXT_MAIN,
                 highlight_words=None):
    """Create a slide with a title and bullet items."""
    set_bg(slide, bg_color)
    add_text(slide, title, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             font_size=40, bold=True, color=title_color)
    y = Inches(1.7)
    for item in items:
        tf = add_text(slide, f"▸  {item}", Inches(0.8), y, Inches(11.5), Inches(0.7),
                       font_size=28, color=item_color)
        y += Inches(0.75)


# ── Build ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# ── SLIDE 1: Title ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "ECC REDMOND  |  MARCH 15, 2026", Inches(1), Inches(1.8),
         Inches(11.3), Inches(0.6), font_size=20, bold=True,
         color=INDIGO_ACC, align=PP_ALIGN.CENTER)
add_text(s, "The Power to Set Free", Inches(1), Inches(2.5),
         Inches(11.3), Inches(1.4), font_size=60, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Sunday School — Middle Schoolers (Grades 6–8)", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, "Romans 8:1-4", Inches(1), Inches(4.9),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)

# ── SLIDE 2: Icebreaker Title ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_TEAL)
add_text(s, "🔓", Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Escape Room Rankings", Inches(1), Inches(2.0),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "You've been locked in a room for an ENTIRE YEAR.\n"
         "No phone. No Wi-Fi. No choosing your own food. Total lockdown.",
         Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.2),
         font_size=26, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "The door opens. You're FREE.", Inches(1), Inches(4.8),
         Inches(11.3), Inches(0.7), font_size=36, bold=True,
         color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s, "What are the FIRST 3 things you do?!", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.7), font_size=32, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "You have 30 seconds. GO!", Inches(1), Inches(6.3),
         Inches(11.3), Inches(0.5), font_size=22, italic=True,
         color=SLATE, align=PP_ALIGN.CENTER)

# ── SLIDE 3: Reveal intro ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x0D, 0x94, 0x88))
add_text(s, "📊", Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s,
         "OK let's see how your answers stack up\nagainst REAL DATA from released prisoners!",
         Inches(1), Inches(3.0), Inches(11.3), Inches(1.6),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDES 4-8: Freedom answers ─────────────────────────────────────────
freedom_data = [
    ("Real Answer #1", "🏠", "Spending time with family"),
    ("Real Answer #2", "🍔", "Eating what you want, where you want"),
    ("Real Answer #3", "🚿", "Taking a long, private shower"),
    ("Real Answer #4", "👕", "Picking out clothes that are YOUR style"),
    ("Real Answer #5", "🌳", "Just being alone outside somewhere"),
]
subtitles = [
    None,
    "First stop = favorite restaurant or a home-cooked meal",
    "No joke — this is a huge deal when you haven't had privacy",
    None,
    None,
]

for (label, emoji, text), sub in zip(freedom_data, subtitles):
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_TEAL)
    add_text(s, label, Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
             font_size=22, bold=True, color=TEAL_ACC, align=PP_ALIGN.CENTER)
    add_text(s, emoji, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
             font_size=72, align=PP_ALIGN.CENTER)
    add_text(s, text, Inches(1), Inches(3.3), Inches(11.3), Inches(1.0),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, sub, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
                 font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# ── SLIDE 9: Icebreaker Transition ──────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "Plot twist...", Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "The Bible says that before we knew Christ,\n"
         "we were ALL prisoners.\n"
         "Not locked in a building —\n"
         "but locked up by sin and death.",
         Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.0),
         font_size=32, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "Last week's sermon was about how\nJesus literally broke us out.",
         Inches(1), Inches(4.8), Inches(11.3), Inches(1.2),
         font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── SLIDE 10: Scripture ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Key Verse", Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"There is therefore now no condemnation for those\n'
         'who are in Christ Jesus. For the law of the Spirit of life\n'
         'has set you free in Christ Jesus from the law of sin and death."',
         Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.5),
         font_size=32, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "— Romans 8:1-2", Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# ── SLIDE 11: Sermon Outline ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "Last Week's Sermon", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=INDIGO_ACC, align=PP_ALIGN.LEFT)
add_text(s, "The Power to Set Free", Inches(1), Inches(1.5), Inches(11.3), Inches(0.9),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
parts = [
    "Part 1:  The Law of Sin and Death Defeated",
    "Part 2:  The Prisoners Freed",
    "Part 3:  The Life in the Spirit",
]
y = Inches(3.0)
for part in parts:
    add_text(s, f"▸  {part}", Inches(1), y, Inches(11), Inches(0.7),
             font_size=32, color=TEXT_MAIN)
    y += Inches(0.85)

# ── SLIDE 12: Key Idea 1 ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE)
add_text(s, "Key Idea 1", Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=LIGHT_BLUE)
add_text(s, "Boss Battle: Christ vs. the Grave Lord", Inches(0.8), Inches(1.2),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    "The final boss — the Grave Lord (sin & death) — was undefeated",
    "Jesus took off His royal armor, disguised Himself as an ordinary soldier",
    "He grabbed the Grave Lord and hurled them both off a cliff into the sea",
    "After 3 days, Jesus rose — because He has the power of life in Him",
    "He returned to the enemy's prison camp and freed every captive",
]
y = Inches(2.5)
for item in items:
    add_text(s, f"▸  {item}", Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.85)

# ── SLIDE 13: Three Victories ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x03, 0x69, 0xA1))
add_text(s, "Three Things Changed for Us", Inches(1), Inches(0.8),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
victories = [
    "⚔️  The war is WON",
    "🔓  The prison doors are OPEN",
    "📜  The contract is CANCELED",
]
y = Inches(2.2)
for v in victories:
    # Draw a rounded-ish box
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(3), y, Inches(7.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.fill.fore_color.brightness = -0.85
    shape.line.color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].text = v
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    y += Inches(1.3)
add_text(s, "We don't owe sin anything anymore.", Inches(1), Inches(6.2),
         Inches(11.3), Inches(0.6), font_size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ── SLIDE 14: Key Idea 2 ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_PURPLE)
add_text(s, "Key Idea 2", Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=PURPLE_ACC)
add_text(s, "The Ultimate Swap", Inches(0.8), Inches(1.2),
         Inches(11.5), Inches(0.9), font_size=46, bold=True, color=WHITE)
items = [
    'From A Tale of Two Cities — Darnay is sentenced to die for his family\'s crimes',
    'His look-alike Carton sneaks into prison, drugs him, swaps clothes',
    'Carton — totally innocent — dies in Darnay\'s place at the guillotine',
    'That\'s what Jesus did: He took our place and gave us His freedom',
]
y = Inches(2.5)
for item in items:
    add_text(s, f"▸  {item}", Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=28, color=TEXT_MAIN)
    y += Inches(0.85)

# ── SLIDE 15: 2 Cor 5:21 ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x58, 0x1C, 0x87))
add_text(s,
         '"For our sake God made him to be sin who knew no sin,\n'
         'so that in him we might become the righteousness of God."',
         Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.5),
         font_size=34, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "— 2 Corinthians 5:21", Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# ── SLIDE 16: Key Idea 3 ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "Key Idea 3", Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=GREEN_ACC)
add_text(s, "Don't Go Back to Jail", Inches(0.8), Inches(1.2),
         Inches(11.5), Inches(0.9), font_size=46, bold=True, color=WHITE)
items = [
    "Jesus didn't break us out so we'd live like escaped criminals",
    "He freed us to live a completely new life in His Spirit",
    "But sometimes we voluntarily walk back into prison...",
    "We start serving false masters that promise freedom but deliver condemnation",
]
y = Inches(2.5)
for item in items:
    add_text(s, f"▸  {item}", Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=28, color=TEXT_MAIN)
    y += Inches(0.85)

# ── SLIDE 17: False Masters ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x15, 0x80, 0x3D))
add_text(s, "False Masters = Terrible Bosses", Inches(0.8), Inches(0.5),
         Inches(11.5), Inches(0.6), font_size=28, bold=True, color=GREEN_ACC)
add_text(s, "They reward you when you win, crush you when you don't",
         Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8),
         font_size=34, bold=True, color=WHITE)
masters = [
    "📱  Social Media — likes = great, left on read = crushed",
    "📊  Grades — A = amazing, bad score = \"I'm a failure\"",
    "👟  Appearance — new outfit = confidence, bad hair day = hide",
    "🏆  Being the Best — winning = freedom, losing = worthless",
]
y = Inches(2.5)
for m in masters:
    add_text(s, f"▸  {m}", Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.85)
add_text(s, "None of these things forgive you when you fail. Only Jesus does.",
         Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
         font_size=24, color=GREEN_ACC)

# ── SLIDE 18: Tim Keller Quote ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s,
         '"Everybody has to live for something. Whatever that something is\n'
         'becomes \'Lord of your life,\' whether you think of it that way or not.\n\n'
         'Jesus is the only Lord who, if you receive him, will fulfill you\n'
         'completely, and, if you fail him, will forgive you eternally."',
         Inches(1.2), Inches(1.5), Inches(10.8), Inches(3.5),
         font_size=30, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "— Tim Keller", Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# ── SLIDE 19: Two Paths ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x29, 0x25, 0x24))
add_text(s, "Which life do you want?", Inches(1), Inches(0.4),
         Inches(11.3), Inches(0.8), font_size=34, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

# Left column: False Master
add_text(s, "Serving a False Master 😩", Inches(0.5), Inches(1.5),
         Inches(5.8), Inches(0.6), font_size=24, bold=True, color=RED_ACC,
         align=PP_ALIGN.CENTER)
false_items = [
    "Can't stop replaying mistakes",
    '"I\'m so stupid" self-talk',
    "Imagining everyone judges you",
    "Judging others by performance",
]
y = Inches(2.3)
for item in false_items:
    add_text(s, f"▸  {item}", Inches(0.8), y, Inches(5.2), Inches(0.6),
             font_size=24, color=LIGHT_RED)
    y += Inches(0.7)

# Right column: Spirit
add_text(s, "Walking in Christ's Spirit 🔥", Inches(6.8), Inches(1.5),
         Inches(5.8), Inches(0.6), font_size=24, bold=True, color=GREEN_ACC,
         align=PP_ALIGN.CENTER)
spirit_items = ["Love", "Joy", "Peace", "Patience",
                "Kindness, Goodness", "Faithfulness", "Gentleness, Self-control"]
y = Inches(2.3)
for item in spirit_items:
    add_text(s, f"▸  {item}", Inches(7.1), y, Inches(5.2), Inches(0.5),
             font_size=22, color=LIGHT_GREEN)
    y += Inches(0.55)

# ── SLIDE 20: Video ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x37, 0x30, 0xA3))
add_text(s, "Recommended Video (8 min)", Inches(1), Inches(1.5),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
add_text(s, "BibleProject: Romans 5–16 Overview", Inches(1), Inches(2.3),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "▶  youtube.com/watch?v=0SVTl4Xa5fY", Inches(1), Inches(3.8),
         Inches(11.3), Inches(0.6), font_size=28, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Click link above or search \"BibleProject Romans 5-16\"\n"
         "Covers Romans 5–16 themes including freedom from condemnation",
         Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# ── SLIDE 21: Key Takeaway ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "THE BOTTOM LINE", Inches(1), Inches(1.0),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"There is therefore now NO condemnation\n'
         'for those who are in Christ Jesus."',
         Inches(1), Inches(2.0), Inches(11.3), Inches(2.0),
         font_size=40, italic=True, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "Not some. Not a little.", Inches(1), Inches(4.5),
         Inches(11.3), Inches(0.6), font_size=32, color=TEXT_MAIN,
         align=PP_ALIGN.CENTER)
add_text(s, "NONE. ZERO. NADA.", Inches(1), Inches(5.3),
         Inches(11.3), Inches(0.8), font_size=46, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

# ── SLIDE 22: Small Group Time ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "💬", Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Small Group Time", Inches(1), Inches(2.8),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Split into groups of 4-5 with a leader", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s,
         "No judgment zone. Share honestly.\nAsk questions. Listen to each other.",
         Inches(1), Inches(5.0), Inches(11.3), Inches(1.0),
         font_size=24, color=INDIGO_ACC, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__),
                   "ECC Sunday School - The Power to Set Free (Slides).pptx")
prs.save(out)
print(f"Saved → {out}")
print(f"Total slides: {len(prs.slides)}")
