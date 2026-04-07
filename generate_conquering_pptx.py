"""Generate PPTX for: ECC Sunday School - Conquering Through Christ (Slides)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Palette --
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
BG_INDIGO    = RGBColor(0x1E, 0x1B, 0x4B)
BG_RED       = RGBColor(0x7F, 0x1D, 0x1D)
BG_RED_MID   = RGBColor(0xB9, 0x1C, 0x1C)
BG_BLUE      = RGBColor(0x0C, 0x4A, 0x6E)
BG_PURPLE    = RGBColor(0x4A, 0x19, 0x42)
BG_GREEN     = RGBColor(0x14, 0x53, 0x2D)
BG_STONE     = RGBColor(0x1C, 0x19, 0x17)

WHITE        = RGBColor(0xF1, 0xF5, 0xF9)
GOLD         = RGBColor(0xFB, 0xBF, 0x24)
LIGHT_BLUE   = RGBColor(0x93, 0xC5, 0xFD)
INDIGO_ACC   = RGBColor(0x81, 0x8C, 0xF8)
TEAL_ACC     = RGBColor(0x5E, 0xEA, 0xD4)
PURPLE_ACC   = RGBColor(0xF0, 0xAB, 0xFC)
GREEN_ACC    = RGBColor(0x86, 0xEF, 0xAC)
RED_ACC      = RGBColor(0xFC, 0xA5, 0xA5)
SLATE        = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MAIN    = RGBColor(0xE2, 0xE8, 0xF0)
CYAN_ACC     = RGBColor(0x7D, 0xD3, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=32, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf


def add_para(tf, text, font_size=32, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    return p


# -- Build --
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# -- SLIDE 1: Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "ECC REDMOND  |  APRIL 5, 2026 \u2014 EASTER SUNDAY", Inches(1), Inches(1.8),
         Inches(11.3), Inches(0.6), font_size=20, bold=True,
         color=INDIGO_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Conquering Through Christ", Inches(1), Inches(2.5),
         Inches(11.3), Inches(1.4), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Sunday School \u2014 Middle Schoolers (Grades 6\u20138)", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, "Romans 8:31-39", Inches(1), Inches(4.9),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)

# -- SLIDE 2: Icebreaker Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_RED)
add_text(s, "\U0001f4aa", Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "UNSTOPPABLE", Inches(1), Inches(2.0),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "I'll describe an obstacle. You shout what can STOP it\n"
         '\u2014 or say "UNSTOPPABLE" if nothing can!',
         Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.4),
         font_size=26, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "Get ready...", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.5), font_size=22, italic=True,
         color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDES 3-8: Obstacles --
obstacles = [
    ("What can stop...", "\U0001f682", "A speeding train?", None),
    ("What can stop...", "\U0001f30a", "A flood of water?", None),
    ("What can stop...", "\U0001f4f1", "A viral TikTok trend?", None),
    ("What can stop...", "\u23f0", "Your alarm clock\non Monday morning?", None),
    ("What can stop...", "\U0001f981", "A lion chasing you?",
     'Pastor Steven mentioned 1 Peter 5:8 \u2014\n"Your adversary, the devil, prowls around like a roaring lion..."'),
]

for (label, emoji, text, sub) in obstacles:
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_RED_MID)
    add_text(s, label, Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
             font_size=22, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
    add_text(s, emoji, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
             font_size=72, align=PP_ALIGN.CENTER)
    add_text(s, text, Inches(1), Inches(3.3), Inches(11.3), Inches(1.2),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, sub, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
                 font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 8: The Big One -- God's Love --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "What can stop...", Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
add_text(s, "\u2764\ufe0f", Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "God's love for you?", Inches(1), Inches(3.3), Inches(11.3), Inches(1.0),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "NOTHING.", Inches(1), Inches(4.8), Inches(11.3), Inches(0.8),
         font_size=48, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 9: Transition \u2014 Easter --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "Happy Easter!", Inches(1), Inches(0.5), Inches(11.3), Inches(1.0),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "Pastor Steven said: God doesn't leave Easter in the past.\n"
         "God wants Easter in our present lives.",
         Inches(1.2), Inches(1.8), Inches(10.8), Inches(1.2),
         font_size=28, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "Christ's death and resurrection aren't just history \u2014\n"
         "they are resurrection power for living TODAY.",
         Inches(1.2), Inches(3.3), Inches(10.8), Inches(1.0),
         font_size=28, color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s,
         "What then shall we say\nto these things?",
         Inches(1), Inches(5.0), Inches(11.3), Inches(1.2),
         font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 10: Scripture Romans 8:31-32 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Key Verse", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"What then shall we say to these things?\n'
         'If God is for us, who is against us?\n'
         'He who did not spare His own Son,\n'
         'but delivered Him over for us all,\n'
         'how will He not also with Him\n'
         'freely give us all things?"',
         Inches(1.2), Inches(1.8), Inches(10.8), Inches(3.5),
         font_size=32, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Romans 8:31-32", Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 11: Scripture Romans 8:37-39 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s,
         '"But in all these things we overwhelmingly\n'
         'conquer through Him who loved us.\n'
         'For I am convinced that neither death, nor life,\n'
         'nor angels, nor principalities, nor things present,\n'
         'nor things to come, nor powers, nor height, nor depth,\n'
         'nor any other created thing will be able to separate us\n'
         'from the love of God that is in Christ Jesus our Lord."',
         Inches(1.2), Inches(1.2), Inches(10.8), Inches(4.5),
         font_size=30, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Romans 8:37-39", Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 12: Sermon Overview --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "Easter Sermon", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=INDIGO_ACC, align=PP_ALIGN.LEFT)
add_text(s, "Conquering Through Christ", Inches(1), Inches(1.5), Inches(11.3), Inches(0.9),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
parts = [
    "The Question:  What shall we say to these things? (v. 31a)",
    "5 Answers:  Five rhetorical questions (vv. 31b-36)",
    "Conclusion:  God's love is the power that overcomes (vv. 37-39)",
]
y = Inches(3.0)
for part in parts:
    add_text(s, "\u25b8  {}".format(part), Inches(1), y, Inches(11), Inches(0.7),
             font_size=30, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 13: Key Idea 1 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE)
add_text(s, "Key Idea 1", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=LIGHT_BLUE)
add_text(s, "God Already Paid the Highest Price", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    "If God gave up His own Son for us, why would He skip the easy stuff?",
    "A man who buys a $500K boat doesn't leave it parked in his driveway",
    "God paid the ultimate price to get you \u2014 He's NOT going to forget about you",
    '"Has God forgotten me?" \u2014 Absolutely not. He\'s right there in your life.',
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 14: Key Idea 2 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x6B, 0x21, 0x63))
add_text(s, "Key Idea 2", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=PURPLE_ACC)
add_text(s, "Nothing Can Condemn You", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    '"Not guilty" is NOT the same as "innocent"',
    "Like a debt that's been paid by someone else \u2014 the debt is GONE",
    "Jesus fully and completely paid our debt of sin. It is finished.",
    "Satan wants you to wallow in guilt. God wants you to confess, get up, and keep walking.",
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 15: Jesus Intercedes --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x44, 0x40, 0x3C))
add_text(s, "Right Now, Jesus Is Praying for You", Inches(1), Inches(0.8),
         Inches(11.3), Inches(0.8), font_size=34, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"Simon, Simon, behold, Satan has demanded\n'
         'to sift you like wheat, but I have prayed\n'
         'for you, that your faith will not fail."',
         Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.5),
         font_size=34, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Luke 22:31-32", Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s,
         "TWO persons of the Trinity \u2014 Jesus AND the Holy Spirit \u2014\n"
         "are praying for you right now.",
         Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.8),
         font_size=24, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 16: Key Idea 3 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "Key Idea 3", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=GREEN_ACC)
add_text(s, 'We Are "Super Conquerors"', Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    '"Overwhelmingly conquer" = super conqueror',
    "A super conqueror takes his enemies and makes them serve him",
    "God takes what Satan uses against us and turns it for our good",
    "We conquer not by our own strength but through Him who loved us",
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 17: Everything Serves God's Purpose --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x03, 0x69, 0xA1))
add_text(s, "All Things Work Together for Good", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.6), font_size=28, bold=True,
         color=CYAN_ACC, align=PP_ALIGN.CENTER)

serve_items = [
    ("\U0001f480", "Death", "frees us from broken bodies \u2192 brings us to Jesus"),
    ("\U0001f331", "Life", "gives us room to grow"),
    ("\U0001f47c", "Angels", "serve us"),
    ("\u2694\ufe0f", "Trials", "strengthen our faith"),
    ("\U0001f30d", "Present things", "teach us only Christ satisfies"),
    ("\U0001f52e", "Future things", "are under God's control"),
]
col1_x = Inches(1.0)
col2_x = Inches(7.0)
y_start = Inches(1.5)
for i, (emoji, title, desc) in enumerate(serve_items):
    if i < 3:
        x = col1_x
        y = y_start + Inches(i * 1.6)
    else:
        x = col2_x
        y = y_start + Inches((i - 3) * 1.6)
    add_text(s, "{} {}".format(emoji, title), x, y, Inches(5.5), Inches(0.7),
             font_size=32, bold=True, color=GOLD)
    add_text(s, desc, x + Inches(0.3), y + Inches(0.65), Inches(5.2), Inches(0.6),
             font_size=22, color=SLATE)

add_text(s, "Everything is made to serve God's saving, sanctifying love work in you.",
         Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 18: 1 John 4:4 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s,
         '"You are from God, little children,\n'
         'and have overcome them;\n'
         'because greater is He who is in you\n'
         'than he who is in the world."',
         Inches(1.2), Inches(1.5), Inches(10.8), Inches(3.5),
         font_size=38, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 1 John 4:4", Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 19: Spurgeon --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x29, 0x25, 0x24))
add_text(s, "C.H. Spurgeon \u2014 A Real Super Conqueror", Inches(0.8), Inches(0.4),
         Inches(11.5), Inches(0.8), font_size=30, bold=True,
         color=GOLD)
items = [
    "One of the 5 greatest preachers in history",
    "At age 20, a prank caused a stampede \u2014 7 died. He suffered lifelong depression",
    "Also suffered severe gout \u2014 bedridden for days at a time",
    "Yet he never gave up. Planted 200+ churches. Started a seminary.",
    "He conquered by holding on to Christ and never letting go",
]
y = Inches(1.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.85)
add_text(s,
         '"You don\'t have to be happy and clappy all the time\nto be faithful to Christ."',
         Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.9),
         font_size=24, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 20: Three Responses --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "How Should We Respond?", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.8), font_size=34, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

responses = [
    "\U0001f305  Bring Easter into your daily life",
    "\U0001f4d6  Love and serve God",
    "\U0001f91d  Care for those God cares for",
]
y = Inches(2.0)
for v in responses:
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(3), y, Inches(7.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.fill.fore_color.brightness = -0.85
    shape.line.color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].text = v
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    y += Inches(1.3)

add_text(s,
         "Difficulties are not just burdens \u2014 they are opportunities to look to God.\n"
         "Easter is resurrection power for living TODAY.",
         Inches(1), Inches(6.0), Inches(11.3), Inches(1.0),
         font_size=22, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 21: Video --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x37, 0x30, 0xA3))
add_text(s, "Recommended Video (11 min)", Inches(1), Inches(1.5),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
add_text(s, "BibleProject: Romans 5\u201316 Overview", Inches(1), Inches(2.3),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "\u25b6  youtube.com/watch?v=0SVTl4Xa5fY", Inches(1), Inches(3.8),
         Inches(11.3), Inches(0.6), font_size=28, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         'Click link above or search "BibleProject Romans 5-16"\n'
         "Covers Romans 5\u201316 including Romans 8 and God's unbreakable love",
         Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 22: Key Takeaway --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "THE BOTTOM LINE", Inches(1), Inches(0.8),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "God already paid the highest price for you.\n"
         "Absolutely NOTHING can separate you from His love.",
         Inches(1), Inches(1.8), Inches(11.3), Inches(1.8),
         font_size=38, italic=True, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "Not your worst day. Not your biggest failure.\nNot your deepest fear.",
         Inches(1), Inches(4.0), Inches(11.3), Inches(1.0),
         font_size=30, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "You are more than a conqueror.", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.8), font_size=44, bold=True,
         color=GREEN_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 23: Small Group Time --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "\U0001f4ac", Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Small Group Time", Inches(1), Inches(2.8),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Split into groups of 4-5 with a leader", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s,
         "No judgment zone. Share honestly.\nAsk questions. Listen to each other.",
         Inches(1), Inches(5.0), Inches(11.3), Inches(1.0),
         font_size=24, color=INDIGO_ACC, align=PP_ALIGN.CENTER)

# -- Save --
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ECC Sunday School - Conquering Through Christ (Slides).pptx")
prs.save(out)
print("Saved -> {}".format(out))
print("Total slides: {}".format(len(prs.slides)))
