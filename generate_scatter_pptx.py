"""Generate PPTX for: ECC Sunday School - Scattered, Not Shattered (Slides)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Palette --
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
BG_INDIGO    = RGBColor(0x1E, 0x1B, 0x4B)
BG_RED       = RGBColor(0x7F, 0x1D, 0x1D)
BG_RED_MID   = RGBColor(0xB9, 0x1C, 0x1C)
BG_BLUE      = RGBColor(0x0C, 0x4A, 0x6E)
BG_PURPLE    = RGBColor(0x4A, 0x19, 0x42)
BG_GREEN     = RGBColor(0x14, 0x53, 0x2D)
BG_STONE     = RGBColor(0x1C, 0x19, 0x17)

WHITE        = RGBColor(0xF1, 0xF5, 0xF9)
GOLD         = RGBColor(0xFB, 0xBF, 0x24)
LIGHT_BLUE   = RGBColor(0x93, 0xC5, 0xFD)
INDIGO_ACC   = RGBColor(0x81, 0x8C, 0xF8)
TEAL_ACC     = RGBColor(0x5E, 0xEA, 0xD4)
PURPLE_ACC   = RGBColor(0xF0, 0xAB, 0xFC)
GREEN_ACC    = RGBColor(0x86, 0xEF, 0xAC)
RED_ACC      = RGBColor(0xFC, 0xA5, 0xA5)
SLATE        = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MAIN    = RGBColor(0xE2, 0xE8, 0xF0)
CYAN_ACC     = RGBColor(0x7D, 0xD3, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=32, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf


def add_para(tf, text, font_size=32, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    return p


# -- Build --
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# -- SLIDE 1: Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "ECC REDMOND  |  APRIL 19, 2026", Inches(1), Inches(1.8),
         Inches(11.3), Inches(0.6), font_size=20, bold=True,
         color=INDIGO_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Scattered, Not Shattered", Inches(1), Inches(2.5),
         Inches(11.3), Inches(1.4), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Sunday School \u2014 Middle Schoolers (Grades 6\u20138)", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, "Acts 8:1-8  |  Pastor Steve Moy", Inches(1), Inches(4.9),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)

# -- SLIDE 2: Icebreaker Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_RED)
add_text(s, "\U0001f4a8", Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "SCATTERED!", Inches(1), Inches(2.0),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "I'll name something that starts together.\n"
         'You shout what could SCATTER it!',
         Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.4),
         font_size=26, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "Get ready...", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.5), font_size=22, italic=True,
         color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDES 3-7: Scatter scenarios --
scenarios = [
    ("What scatters...", "\U0001f342", "A pile of leaves\non the ground?", None),
    ("What scatters...", "\U0001f426", "A flock of pigeons\nin the park?", None),
    ("What scatters...", "\U0001f46b", "Your friend group\nat school?", None),
    ("What scatters...", "\U0001f52e", "A bag of marbles\ndropped on the floor?", None),
    ("What scatters...", "\U0001f33c", "Seeds in a\ndandelion puff?",
     "And they grow into NEW flowers wherever they land..."),
]

for (label, emoji, text, sub) in scenarios:
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_RED_MID)
    add_text(s, label, Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
             font_size=22, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
    add_text(s, emoji, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
             font_size=72, align=PP_ALIGN.CENTER)
    add_text(s, text, Inches(1), Inches(3.3), Inches(11.3), Inches(1.2),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, sub, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
                 font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 8: The Big One -- Early Church --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "What scattered...", Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
add_text(s, "\u26ea", Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "The early church\nin Jerusalem?", Inches(1), Inches(3.3), Inches(11.3), Inches(1.0),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "PERSECUTION.", Inches(1), Inches(4.8), Inches(11.3), Inches(0.8),
         font_size=48, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
add_text(s, "But like dandelion seeds, they grew new churches\neverywhere they landed.",
         Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.8),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 9: Transition -- Star Wars --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "\U0001f30c", Inches(1), Inches(0.3), Inches(11.3), Inches(1.0),
         font_size=60, align=PP_ALIGN.CENTER)
add_text(s, "Disruption \u2192 Bigger Purpose", Inches(1), Inches(1.2),
         Inches(11.3), Inches(0.8), font_size=40, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "Pastor Steve's Star Wars comparison:\n"
         "Luke Skywalker was comfortable on Tatooine \u2014\n"
         "harvesting water, minding his own business.",
         Inches(1.2), Inches(2.4), Inches(10.8), Inches(1.5),
         font_size=26, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "Then two droids showed up. His family died.\n"
         "His home was destroyed. That painful disruption\n"
         "pushed him into something greater \u2014 the Rebellion.",
         Inches(1.2), Inches(4.0), Inches(10.8), Inches(1.5),
         font_size=26, color=TEAL_ACC, align=PP_ALIGN.CENTER)
add_text(s,
         "The early church didn't choose to leave Jerusalem \u2014\n"
         "but God used that disruption to scatter them\n"
         "exactly where Jesus told them to go.",
         Inches(1.2), Inches(5.8), Inches(10.8), Inches(1.0),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 10: Scripture Acts 8:1 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Acts 8:1", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"And Saul approved of his execution.\n'
         'And there arose on that day\n'
         'a great persecution against the church\n'
         'in Jerusalem, and they were all scattered\n'
         'throughout the regions of Judea and Samaria,\n'
         'except the apostles."',
         Inches(1.2), Inches(1.8), Inches(10.8), Inches(3.5),
         font_size=32, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Acts 8:1 (ESV)", Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, 'Greek "approved" = suneodokeo \u2014 "to be pleased with"',
         Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
         font_size=20, italic=True, color=RED_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 11: Scripture Acts 8:4-5 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Acts 8:4-5", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"Now those who were scattered\n'
         'went about preaching the word.\n'
         'Philip went down to the city of Samaria\n'
         'and proclaimed to them the Christ."',
         Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.0),
         font_size=34, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Acts 8:4-5 (ESV)", Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "They didn't scatter and hide. They scattered and preached.",
         Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
         font_size=22, color=TEAL_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 12: Scripture Acts 8:8 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Acts 8:8", Inches(1), Inches(1.2), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
add_text(s,
         '"So there was much joy\nin that city."',
         Inches(1.2), Inches(2.5), Inches(10.8), Inches(2.0),
         font_size=44, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 Acts 8:8 (ESV)", Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s,
         "Persecution  \u2192  Scattering  \u2192  Preaching  \u2192  Healing  \u2192  JOY",
         Inches(1), Inches(6.0), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 13: Sermon Overview --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "Pastor Steve Moy", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=INDIGO_ACC, align=PP_ALIGN.LEFT)
add_text(s, "Scattered, Not Shattered", Inches(1), Inches(1.5), Inches(11.3), Inches(0.9),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
parts = [
    "1)  God Uses Disruption to Scatter His People",
    "2)  The Courage of the Scattered",
    "3)  The Joy of the Unexpected Harvest",
]
y = Inches(3.0)
for part in parts:
    add_text(s, "\u25b8  {}".format(part), Inches(1), y, Inches(11), Inches(0.7),
             font_size=30, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 14: Key Idea 1 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE)
add_text(s, "Key Idea 1", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=LIGHT_BLUE)
add_text(s, "God Uses Disruption to Move His People", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    "Jesus said go to all Judea and Samaria (Acts 1:8) \u2014 but the church stayed in Jerusalem",
    "They were comfortable \u2014 gathering, eating, worshiping, growing together",
    "Persecution forced them outward \u2014 exactly where Jesus told them to go",
    "God uses disruption not to punish us but to propel us",
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 15: Key Idea 2 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x6B, 0x21, 0x63))
add_text(s, "Key Idea 2", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=PURPLE_ACC)
add_text(s, "Ordinary People, Extraordinary Mission", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    "The apostles stayed in Jerusalem \u2014 ordinary believers scattered and preached",
    'Philip was a deacon from Acts 6 \u2014 "in essence, he was a waiter"',
    "He went to Samaria \u2014 Jews and Samaritans had hundreds of years of hostility",
    'Mission belongs to all of us \u2014 not just pastors and missionaries',
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 16: 2 Timothy 1:7 --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s,
         '"For God gave us a spirit\n'
         'not of fear but of power\n'
         'and of love and self-discipline."',
         Inches(1.2), Inches(1.5), Inches(10.8), Inches(3.0),
         font_size=38, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "\u2014 2 Timothy 1:7", Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s,
         "This is the power that compelled ordinary believers \u2014\n"
         "not trained preachers \u2014 to share the gospel\n"
         "even while being hunted.",
         Inches(1.5), Inches(5.8), Inches(10.3), Inches(1.0),
         font_size=24, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 17: Tim Keller Quote --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x29, 0x25, 0x24))
add_text(s, "Tim Keller", Inches(1), Inches(1.2), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s,
         '"The Gospel is this;\n'
         'We are more sinful and flawed\n'
         'than we ever dared to believe,\n'
         'yet more loved and accepted\n'
         'than we ever dared to hope."',
         Inches(1.2), Inches(2.2), Inches(10.8), Inches(3.5),
         font_size=36, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "This gospel is not reserved for a few \u2014\n"
         "it's entrusted to all of us.",
         Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.8),
         font_size=24, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 18: COVID Parallel --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x03, 0x69, 0xA1))
add_text(s, "Modern-Day Scattering", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=CYAN_ACC)
add_text(s, "The COVID Pandemic (2020-2022)", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=38, bold=True, color=WHITE)
items = [
    "Churches were forcibly closed \u2014 believers scattered to living rooms and Zoom",
    "Surge in digital missions and local neighborhood care",
    "The gospel was taken out of churches and into the streets",
    "Hundreds still tune in to ECC's messages online because of that scattering",
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)
add_text(s, "Acts 8 isn't just ancient history \u2014 it still happens today.",
         Inches(1), Inches(6.0), Inches(11.3), Inches(0.6),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 19: Key Idea 3 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "Key Idea 3", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=GREEN_ACC)
add_text(s, "Scattering Leads to Unexpected Joy", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=42, bold=True, color=WHITE)
items = [
    'Crowds paid attention to Philip with "one accord"',
    "Unclean spirits cast out, paralyzed and lame healed",
    '"So there was much joy in that city" (v. 8)',
    "The gospel doesn't just fix problems \u2014 it restores joy to communities",
]
y = Inches(2.5)
for item in items:
    add_text(s, "\u25b8  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 20: The Journey Chain --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x03, 0x69, 0xA1))
add_text(s, "From Scattering to Joy", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.6), font_size=28, bold=True,
         color=CYAN_ACC, align=PP_ALIGN.CENTER)

chain_top = [
    ("Stephen\nmartyred", GOLD),
    ("\u2192", SLATE),
    ("Persecution\nrises", GOLD),
    ("\u2192", SLATE),
    ("Believers\nscattered", GOLD),
]
x = Inches(0.5)
for text, color in chain_top:
    if text == "\u2192":
        add_text(s, text, x, Inches(2.0), Inches(0.8), Inches(0.8),
                 font_size=36, color=color, align=PP_ALIGN.CENTER)
        x += Inches(0.8)
    else:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, Inches(1.5), Inches(2.8), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
        shape.fill.fore_color.brightness = -0.85
        shape.line.color.rgb = RGBColor(0xFB, 0xBF, 0x24)
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = GOLD
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += Inches(2.8)

chain_bot = [
    ("They preach\neverywhere", GOLD),
    ("\u2192", SLATE),
    ("Philip goes\nto Samaria", GOLD),
    ("\u2192", SLATE),
    ("MUCH\nJOY!", GREEN_ACC),
]
x = Inches(0.5)
for text, color in chain_bot:
    if text == "\u2192":
        add_text(s, text, x, Inches(4.2), Inches(0.8), Inches(0.8),
                 font_size=36, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(0.8)
    else:
        is_joy = (text == "MUCH\nJOY!")
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, Inches(3.7), Inches(2.8), Inches(1.0))
        shape.fill.solid()
        if is_joy:
            shape.fill.fore_color.rgb = RGBColor(0x86, 0xEF, 0xAC)
            shape.fill.fore_color.brightness = -0.85
            shape.line.color.rgb = GREEN_ACC
        else:
            shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
            shape.fill.fore_color.brightness = -0.85
            shape.line.color.rgb = GOLD
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += Inches(2.8)

add_text(s,
         "God turned the church's worst day\ninto the gospel's greatest advance.",
         Inches(1), Inches(5.8), Inches(11.3), Inches(1.0),
         font_size=24, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 21: C.S. Lewis Quote --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "C.S. Lewis", Inches(1), Inches(1.5), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s,
         '"Joy is the serious\nbusiness of heaven."',
         Inches(1.2), Inches(2.5), Inches(10.8), Inches(2.5),
         font_size=48, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "The joy in Samaria wasn't quiet or private \u2014\n"
         "it was overflowing, communal, and visible.",
         Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.0),
         font_size=24, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 22: Application --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "How Should We Respond?", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.8), font_size=34, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

responses = [
    "\U0001f309  Equip EVERYONE to share the Gospel \u2014 be bridge builders",
    "\U0001f3af  See yourself on mission \u2014 presence of peace, integrity, grace",
    "\U0001f4aa  Live with courage \u2014 build relationships beyond your comfort zone",
]
y = Inches(2.0)
for v in responses:
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(2), y, Inches(9.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.fill.fore_color.brightness = -0.85
    shape.line.color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].text = v
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    y += Inches(1.3)

add_text(s,
         'Philip was a table server. God used him to bring joy to a whole city.\n'
         'God can use YOU too.',
         Inches(1), Inches(6.0), Inches(11.3), Inches(1.0),
         font_size=22, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 23: Video --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x37, 0x30, 0xA3))
add_text(s, "Recommended Video (8 min)", Inches(1), Inches(1.5),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
add_text(s, "BibleProject: Acts 8\u201312 Overview", Inches(1), Inches(2.3),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "\u25b6  youtube.com/watch?v=Z-17KxKjvfM", Inches(1), Inches(3.8),
         Inches(11.3), Inches(0.6), font_size=28, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         'Click link above or search "BibleProject Acts 8-12"\n'
         "Covers Acts 8-12 including the scattering and Philip's ministry in Samaria",
         Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 24: Key Takeaway --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x4C, 0x1D, 0x95))
add_text(s, "THE BOTTOM LINE", Inches(1), Inches(0.8),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "God turns our scattering into His strategy.\n"
         "Every believer \u2014 not just pastors \u2014 is called\n"
         "to share Christ wherever they are scattered.",
         Inches(1), Inches(1.8), Inches(11.3), Inches(1.8),
         font_size=36, italic=True, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         'Philip was "in essence, a waiter."\nGod used him to change a city.',
         Inches(1), Inches(4.0), Inches(11.3), Inches(1.0),
         font_size=30, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "Scattered, but never shattered.", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.8), font_size=44, bold=True,
         color=GREEN_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 25: Small Group Time --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "\U0001f4ac", Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Small Group Time", Inches(1), Inches(2.8),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Split into groups of 4-5 with a leader", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s,
         "No judgment zone. Share honestly.\nAsk questions. Listen to each other.",
         Inches(1), Inches(5.0), Inches(11.3), Inches(1.0),
         font_size=24, color=INDIGO_ACC, align=PP_ALIGN.CENTER)

# -- Save --
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ECC Sunday School - The Power to Scatter (Slides).pptx")
prs.save(out)
print("Saved -> {}".format(out))
print("Total slides: {}".format(len(prs.slides)))
