"""Generate PPTX for: ECC Sunday School - Your Mission in the Great Commission (Slides)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Palette --
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
BG_INDIGO    = RGBColor(0x1E, 0x1B, 0x4B)
BG_AMBER     = RGBColor(0x42, 0x20, 0x06)
BG_AMBER_MID = RGBColor(0xB4, 0x53, 0x09)
BG_BLUE      = RGBColor(0x0C, 0x4A, 0x6E)
BG_BLUE_MID  = RGBColor(0x03, 0x69, 0xA1)
BG_PURPLE    = RGBColor(0x4A, 0x19, 0x42)
BG_PURPLE_MID= RGBColor(0x6B, 0x21, 0x63)
BG_VIOLET    = RGBColor(0x4C, 0x1D, 0x95)
BG_GREEN     = RGBColor(0x14, 0x53, 0x2D)
BG_STONE     = RGBColor(0x1C, 0x19, 0x17)
BG_STONE_MID = RGBColor(0x29, 0x25, 0x24)

WHITE        = RGBColor(0xF1, 0xF5, 0xF9)
GOLD         = RGBColor(0xFB, 0xBF, 0x24)
AMBER        = RGBColor(0xFC, 0xD3, 0x4D)
LIGHT_BLUE   = RGBColor(0x93, 0xC5, 0xFD)
INDIGO_ACC   = RGBColor(0x81, 0x8C, 0xF8)
TEAL_ACC     = RGBColor(0x5E, 0xEA, 0xD4)
PURPLE_ACC   = RGBColor(0xF0, 0xAB, 0xFC)
GREEN_ACC    = RGBColor(0x86, 0xEF, 0xAC)
RED_ACC      = RGBColor(0xFC, 0xA5, 0xA5)
SLATE        = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MAIN    = RGBColor(0xE2, 0xE8, 0xF0)
CYAN_ACC     = RGBColor(0x7D, 0xD3, 0xFC)
PURPLE_LIGHT = RGBColor(0xC4, 0xB5, 0xFD)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=32, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf


def add_card(slide, x, y, w, h, fill_color, border_color, text,
             font_size=22, font_color=GOLD, bold=True, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.fill.fore_color.brightness = -0.85
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.color.rgb = font_color
    tf.paragraphs[0].alignment = align
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


# -- Build --
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# -- SLIDE 1: Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "ECC REDMOND  |  APRIL 26, 2026", Inches(1), Inches(1.8),
         Inches(11.3), Inches(0.6), font_size=20, bold=True,
         color=INDIGO_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Your Mission in the Great Commission", Inches(1), Inches(2.5),
         Inches(11.3), Inches(1.4), font_size=52, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Sunday School — Middle Schoolers (Grades 6–8)", Inches(1), Inches(4.4),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, "Acts 13:1-5  |  Elder Elton Lee", Inches(1), Inches(5.1),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)

# -- SLIDE 2: Icebreaker Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_AMBER)
add_text(s, "\U0001f3ac", Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "MISSION: POSSIBLE!", Inches(1), Inches(2.0),
         Inches(11.3), Inches(1.2), font_size=52, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "Every great hero gets a mission.\n"
         "Let’s see if you can name them.",
         Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.4),
         font_size=26, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "Round 1 — Who’s the mission?", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.5), font_size=22, italic=True,
         color=AMBER, align=PP_ALIGN.CENTER)

# -- SLIDES 3-7: Round 1 missions --
missions_round1 = [
    ("\U0001f9b8", '"With great power comes\nresponsibility.\nMission: protect New York."', "Spider-Man"),
    ("\U0001f680", '"That’s one small step for man…\nMission: walk on the moon."', "Apollo 11 / Neil Armstrong"),
    ("\U0001f420", '"Cross the entire ocean to find\none little fish.\nMission: rescue your son."', "Marlin (Finding Nemo)"),
    ("⚡", '"Take a magic ring to a volcano\nand throw it in.\nMission: destroy the One Ring."', "Frodo / The Fellowship"),
]

for emoji, text, answer in missions_round1:
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_AMBER_MID)
    add_text(s, "Whose mission?", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
             font_size=22, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_text(s, emoji, Inches(1), Inches(1.6), Inches(11.3), Inches(1.2),
             font_size=72, align=PP_ALIGN.CENTER)
    add_text(s, text, Inches(1), Inches(3.0), Inches(11.3), Inches(2.0),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "→  {}".format(answer), Inches(1), Inches(5.7),
             Inches(11.3), Inches(0.6), font_size=24, italic=True,
             color=AMBER, align=PP_ALIGN.CENTER)

# -- SLIDE 7: The Big One --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "Whose mission?", Inches(1), Inches(0.6), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
add_text(s, "\U0001f31f", Inches(1), Inches(1.4), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s,
         '"Go therefore and make disciples\nof all nations…"\nMission:  ?',
         Inches(1), Inches(2.8), Inches(11.3), Inches(2.0),
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "EVERY CHRISTIAN.", Inches(1), Inches(5.0), Inches(11.3), Inches(0.8),
         font_size=44, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
add_text(s, "Including you. Yes — even at age 12.",
         Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 8: Round 2 Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_AMBER)
add_text(s, "\U0001f4dc", Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, 'Round 2: How does the hero\nget the mission?',
         Inches(1), Inches(2.3), Inches(11.3), Inches(2.2),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Quick — match each one!", Inches(1), Inches(5.5),
         Inches(11.3), Inches(0.6), font_size=22, italic=True,
         color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 9: Receivers grid --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "How the Mission Arrives", Inches(1), Inches(0.5),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

receivers = [
    ("\U0001f3ac", "Mission: Impossible", "self-destructing tape"),
    ("\U0001f3f0", "Frodo", "Council of Elrond"),
    ("\U0001f30c", "Star Wars", "hologram in R2-D2"),
    ("\U0001f4dc", "Avengers", "Nick Fury at your door"),
]
y = Inches(1.7)
for emoji, name, how in receivers:
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1.5), y, Inches(10.3), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.fill.fore_color.brightness = -0.95
    shape.line.color.rgb = SLATE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = "{}  {}  —  {}".format(emoji, name, how)
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = TEXT_MAIN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    y += Inches(0.85)

add_text(s, "And in Acts 13?", Inches(1), Inches(5.6),
         Inches(11.3), Inches(0.6), font_size=28, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "The Holy Spirit speaks while the church is worshiping and fasting.",
         Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.6),
         font_size=24, color=TEAL_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 10: Transition --
s = prs.slides.add_slide(blank)
set_bg(s, BG_VIOLET)
add_text(s, "\U0001f3a7", Inches(1), Inches(0.8), Inches(11.3), Inches(1.0),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s,
         "If God wanted to give YOU\na mission this week…",
         Inches(1), Inches(2.0), Inches(11.3), Inches(2.0),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "… would you actually be paying enough attention\nto hear Him?",
         Inches(1), Inches(4.2), Inches(11.3), Inches(1.5),
         font_size=30, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "That’s what today’s lesson is about.",
         Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 11: Scripture Acts 13:1-2 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Acts 13:1-2", Inches(1), Inches(0.6), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"Now in the church at Antioch\nthere were prophets and teachers…\n'
         'While they were worshiping the Lord and fasting,\n'
         "the Holy Spirit said,\n"
         '‘Set apart for me Barnabas and Saul\nfor the work to which I have called them.’"',
         Inches(1), Inches(1.6), Inches(11.3), Inches(4.5),
         font_size=26, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "— Acts 13:1-2 (NIV)", Inches(1), Inches(6.4), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 12: Scripture Acts 13:3-5 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Acts 13:3-5", Inches(1), Inches(0.6), Inches(11.3), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         '"So after they had fasted and prayed,\n'
         "they placed their hands on them and sent them off.\n"
         "The two of them, sent on their way by the Holy Spirit,\n"
         "went down to Seleucia and sailed from there to Cyprus.\n"
         'When they arrived at Salamis, they proclaimed the word of God…"',
         Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5),
         font_size=24, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "— Acts 13:3-5 (NIV)", Inches(1), Inches(6.4), Inches(11.3), Inches(0.6),
         font_size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 13: Sermon Overview --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "Elder Elton Lee", Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=INDIGO_ACC, align=PP_ALIGN.LEFT)
add_text(s, "Your Mission in the Great Commission", Inches(1), Inches(1.5),
         Inches(11.3), Inches(0.9), font_size=38, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
parts = [
    "1)  The Holy Spirit Initiates God’s Global Mission",
    "2)  The Church Responds — as Sender AND Sent",
    "3)  Your Mission: Learn → Discern → Confirm",
]
y = Inches(3.0)
for part in parts:
    add_text(s, "▸  {}".format(part), Inches(1), y, Inches(11.3), Inches(0.7),
             font_size=28, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 14: From Acts 8 to Acts 13 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "Quick Recap from Last Week", Inches(0.8), Inches(0.4),
         Inches(11.5), Inches(0.6), font_size=26, bold=True, color=GREEN_ACC)
add_text(s, "From Scattered → Sending", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=38, bold=True, color=WHITE)
add_card(s, Inches(0.8), Inches(2.4), Inches(5.7), Inches(1.0),
         GOLD, GOLD,
         "Acts 8:\nPersecution scatters believers",
         font_size=20, font_color=GOLD)
add_text(s, "→", Inches(6.55), Inches(2.6), Inches(0.6), Inches(0.6),
         font_size=36, color=SLATE, align=PP_ALIGN.CENTER)
add_card(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(1.0),
         GOLD, GOLD,
         "They plant the Antioch church",
         font_size=22, font_color=GOLD)
add_card(s, Inches(0.8), Inches(3.8), Inches(6.0), Inches(1.0),
         GOLD, GOLD,
         "Acts 13:\nAntioch sends Paul + Barnabas",
         font_size=20, font_color=GOLD)
add_text(s, "→", Inches(6.85), Inches(4.0), Inches(0.6), Inches(0.6),
         font_size=36, color=SLATE, align=PP_ALIGN.CENTER)
add_card(s, Inches(7.5), Inches(3.8), Inches(5.0), Inches(1.0),
         GREEN_ACC, GREEN_ACC,
         "Gospel goes to the\nends of the earth",
         font_size=20, font_color=GREEN_ACC)
add_text(s,
         "The persecuted became the senders.\nGod writes the best comeback stories.",
         Inches(1), Inches(5.6), Inches(11.3), Inches(1.2),
         font_size=24, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 15: Key Idea 1 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE)
add_text(s, "Key Idea 1", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=LIGHT_BLUE)
add_text(s, "The Holy Spirit Initiates God’s Mission", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=40, bold=True, color=WHITE)
items = [
    "The church wasn’t in a strategy meeting — they were worshiping and fasting",
    "The Holy Spirit’s fingerprints are everywhere in Acts (13:4, 13:9, 16:6-9)",
    "God already has a plan — we just need to find our place in it",
    '"God is always calling — are we responding?"',
]
y = Inches(2.5)
for item in items:
    add_text(s, "▸  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 16: Listening + Responding --
s = prs.slides.add_slide(blank)
set_bg(s, BG_PURPLE_MID)
add_text(s, "Two-Part Test", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=PURPLE_ACC)
add_text(s,
         "Are You Actually Listening?\nAre You Actually Responding?",
         Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.8),
         font_size=38, bold=True, color=WHITE)
items = [
    "Listening: put down the 5x2 inch screen long enough to hear",
    "Responding: don’t ghost God when His call is uncomfortable",
    '"Sometimes His calling is costly. That’s the point of being a living sacrifice."',
]
y = Inches(3.4)
for item in items:
    add_text(s, "▸  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 17: Key Idea 2 Title --
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE_MID)
add_text(s, "Key Idea 2", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=CYAN_ACC)
add_text(s, "The Church Responds:\nSender AND Sent",
         Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5),
         font_size=44, bold=True, color=WHITE)
add_text(s,
         "The Holy Spirit initiates.\nThe Church responds — in two roles.",
         Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5),
         font_size=28, color=TEXT_MAIN)

# -- SLIDE 18: Sender --
s = prs.slides.add_slide(blank)
set_bg(s, BG_GREEN)
add_text(s, "The Church as SENDER", Inches(0.8), Inches(0.4),
         Inches(11.5), Inches(0.6), font_size=28, bold=True, color=GREEN_ACC)
add_text(s, "Antioch gave THREE things:", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=38, bold=True, color=WHITE)
items = [
    "\U0001f4b0  Resources — Antioch was wealthy. They gave generously.",
    "\U0001f494  Their best people — Barnabas + Paul, no leftovers",
    "\U0001f64f  Prayer + encouragement — they backed their missionaries",
]
y = Inches(2.5)
for item in items:
    add_text(s, "▸  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.8),
             font_size=26, color=TEXT_MAIN)
    y += Inches(0.95)

add_text(s,
         "ECC parallel: ~$500K/yr, 40+ long-term missionaries, 17 countries, 7 from ECC.\nMission Fund needs YOUR earmarked giving.",
         Inches(1), Inches(5.7), Inches(11.3), Inches(1.2),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 19: Sent --
s = prs.slides.add_slide(blank)
set_bg(s, BG_PURPLE_MID)
add_text(s, "The Church as SENT", Inches(0.8), Inches(0.4),
         Inches(11.5), Inches(0.6), font_size=28, bold=True, color=PURPLE_ACC)
add_text(s, 'Attributes of "the Sent":', Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=38, bold=True, color=WHITE)

attrs = [
    ("\U0001f3af", "Missional", "share with anyone"),
    ("\U0001f4c5", "Available", "God’s agenda"),
    ("\U0001f504", "Flexible", "drop your plans"),
    ("✋", "Willing", "say yes"),
]
positions = [
    (Inches(1.0), Inches(2.3)),
    (Inches(7.0), Inches(2.3)),
    (Inches(1.0), Inches(3.5)),
    (Inches(7.0), Inches(3.5)),
]
for (emoji, name, desc), (x, y) in zip(attrs, positions):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x, y, Inches(5.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.fill.fore_color.brightness = -0.93
    shape.line.color.rgb = PURPLE_ACC
    shape.line.width = Pt(1)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "{}  {} — {}".format(emoji, name, desc)
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_MAIN
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)

# Servant heart - full width below
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.fill.fore_color.brightness = -0.93
shape.line.color.rgb = PURPLE_ACC
shape.line.width = Pt(1)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "\U0001f49d  Servant heart — selfless (like Paul: beaten, jailed, faithful)"
p.font.size = Pt(24)
p.font.color.rgb = TEXT_MAIN
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# -- SLIDE 20: Plot Twist Paul --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE)
add_text(s, "Plot Twist", Inches(1), Inches(0.6), Inches(11.3), Inches(0.6),
         font_size=26, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_text(s, "Paul used to be a persecutor.",
         Inches(1), Inches(1.4), Inches(11.3), Inches(1.0),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
         "He was part of the very persecution that scattered\n"
         "believers to Antioch — the church that’s now\n"
         "sending him out.",
         Inches(1), Inches(2.6), Inches(11.3), Inches(1.7),
         font_size=24, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         '"He is a chosen instrument of mine\n'
         'to carry my name before the Gentiles and kings…"',
         Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.4),
         font_size=28, italic=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
add_text(s, "— Acts 9:15", Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
         font_size=20, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "If God can use Paul, He can use you.",
         Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
         font_size=22, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 21: Philippians 1:21 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_VIOLET)
add_text(s,
         '"For to me, to live is Christ\nand to die is gain."',
         Inches(1), Inches(2.0), Inches(11.3), Inches(2.5),
         font_size=46, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "— Philippians 1:21", Inches(1), Inches(4.8),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, "Paul fully embodied the surrendered life.",
         Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
         font_size=24, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 22: Key Idea 3 --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "Key Idea 3", Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
         font_size=28, bold=True, color=INDIGO_ACC)
add_text(s, "Your Mission in the Great Commission", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=38, bold=True, color=WHITE)
items = [
    "All are called to the Great Commission",
    "Not all are called to cross-cultural missions",
    "Start local — but don’t stop there",
    '"The Great Commission is not the Great Suggestion." — Rick Warren',
]
y = Inches(2.5)
for item in items:
    add_text(s, "▸  {}".format(item), Inches(0.8), y, Inches(11.5), Inches(0.7),
             font_size=24, color=TEXT_MAIN)
    y += Inches(0.85)

# -- SLIDE 23: Three Steps Card --
s = prs.slides.add_slide(blank)
set_bg(s, BG_BLUE_MID)
add_text(s, "Demystifying \"Calling\"", Inches(0.8), Inches(0.4),
         Inches(11.5), Inches(0.6), font_size=26, bold=True, color=CYAN_ACC)
add_text(s, "Three Steps to Find Yours", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=36, bold=True, color=WHITE)

steps = [
    ("STEP 1", "\U0001f4d6  LEARN", "Posture of a student. Explore what God is doing.\nTake a class, browse joshuaproject.net, attend a missions conference."),
    ("STEP 2", "\U0001f9ed  DISCERN", "Look at: Passion · Gifting · Opportunity · Prayer · Wise counsel.\nWhat’s God wired you for?"),
    ("STEP 3", "✅  CONFIRM", "Try it. Watch for the Holy Spirit and others to confirm.\nPastor Steve Moy was already pastoring before he was ordained."),
]
y = Inches(2.2)
for step_num, title, body in steps:
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.8), y, Inches(11.7), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.fill.fore_color.brightness = -0.95
    shape.line.color.rgb = SLATE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = step_num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p3 = tf.add_paragraph()
    p3.text = body
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MAIN
    y += Inches(1.65)

# -- SLIDE 24: 5 Discern Questions --
s = prs.slides.add_slide(blank)
set_bg(s, BG_PURPLE_MID)
add_text(s, "5 Discern Questions", Inches(0.8), Inches(0.4),
         Inches(11.5), Inches(0.6), font_size=28, bold=True, color=PURPLE_ACC)
add_text(s, "Ask Yourself…", Inches(0.8), Inches(1.0),
         Inches(11.5), Inches(0.9), font_size=38, bold=True, color=WHITE)

questions = [
    ("\U0001f525", "Passion", "Do I care?"),
    ("\U0001f381", "Gifting", "Am I wired for it?"),
    ("\U0001f6aa", "Opportunity", "Is a door opening?"),
    ("\U0001f64f", "Prayer + Spirit", "Am I asking?"),
]
positions = [
    (Inches(1.0), Inches(2.3)),
    (Inches(7.0), Inches(2.3)),
    (Inches(1.0), Inches(3.5)),
    (Inches(7.0), Inches(3.5)),
]
for (emoji, name, desc), (x, y) in zip(questions, positions):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x, y, Inches(5.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.fill.fore_color.brightness = -0.93
    shape.line.color.rgb = PURPLE_ACC
    shape.line.width = Pt(1)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "{}  {} — {}".format(emoji, name, desc)
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_MAIN
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)

shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.fill.fore_color.brightness = -0.93
shape.line.color.rgb = PURPLE_ACC
shape.line.width = Pt(1)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "\U0001f474  Wise Counsel — What do mature believers see in me?"
p.font.size = Pt(24)
p.font.color.rgb = TEXT_MAIN
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# -- SLIDE 25: Rick Warren Quote --
s = prs.slides.add_slide(blank)
set_bg(s, BG_STONE_MID)
add_text(s, "Pastor Rick Warren", Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
         font_size=24, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s,
         '"A great commitment to the\n'
         "Great Commandment\n"
         "and the Great Commission\n"
         'will make you a great Christian."',
         Inches(1), Inches(2.0), Inches(11.3), Inches(3.5),
         font_size=36, italic=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s,
         "Love God + love others + go and make disciples\n= a life that matters.",
         Inches(1), Inches(5.8), Inches(11.3), Inches(1.0),
         font_size=24, color=GOLD, align=PP_ALIGN.CENTER)

# -- SLIDE 26: Video --
s = prs.slides.add_slide(blank)
set_bg(s, RGBColor(0x37, 0x30, 0xA3))
add_text(s, "Recommended Video (~10 min)", Inches(1), Inches(1.5),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "BibleProject: Acts 13–28 Overview", Inches(1), Inches(2.3),
         Inches(11.3), Inches(0.8), font_size=36, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "▶  youtube.com/watch?v=Z6Yj5R3pqM4", Inches(1), Inches(3.8),
         Inches(11.3), Inches(0.6), font_size=28, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         'Click link above or search "BibleProject Acts 13-28"\n'
         "Covers Paul’s missionary journeys, starting from Acts 13",
         Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
         font_size=22, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# -- SLIDE 27: Key Takeaway --
s = prs.slides.add_slide(blank)
set_bg(s, BG_VIOLET)
add_text(s, "THE BOTTOM LINE", Inches(1), Inches(0.6),
         Inches(11.3), Inches(0.6), font_size=24, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "The Holy Spirit initiates.\n"
         "The Church responds — as Sender and as Sent.\n\n"
         "Your mission is real.\n"
         "Find it through Learn → Discern → Confirm.",
         Inches(1), Inches(1.6), Inches(11.3), Inches(3.5),
         font_size=32, italic=True, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "God already has a plan.",
         Inches(1), Inches(5.4), Inches(11.3), Inches(0.7),
         font_size=30, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
add_text(s, "You just need to find your place in it.",
         Inches(1), Inches(6.2), Inches(11.3), Inches(0.8),
         font_size=36, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)

# -- SLIDE 28: Small Group Time --
s = prs.slides.add_slide(blank)
set_bg(s, BG_INDIGO)
add_text(s, "\U0001f4ac", Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
         font_size=72, align=PP_ALIGN.CENTER)
add_text(s, "Small Group Time", Inches(1), Inches(2.8),
         Inches(11.3), Inches(1.2), font_size=56, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Split into groups of 4-5 with a leader", Inches(1), Inches(4.2),
         Inches(11.3), Inches(0.6), font_size=26, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s,
         "No judgment zone. Share honestly.\nAsk questions. Listen to each other.",
         Inches(1), Inches(5.0), Inches(11.3), Inches(1.0),
         font_size=24, color=INDIGO_ACC, align=PP_ALIGN.CENTER)

# -- Save --
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ECC Sunday School - Your Mission in the Great Commission (Slides).pptx")
prs.save(out)
print("Saved -> {}".format(out))
print("Total slides: {}".format(len(prs.slides)))
